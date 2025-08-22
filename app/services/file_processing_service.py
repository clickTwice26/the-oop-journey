import os
import PyPDF2
from PIL import Image
import docx
from pptx import Presentation
import base64
from io import BytesIO

class FileProcessingService:
    """Service for processing uploaded files and extracting text content"""
    
    @staticmethod
    def extract_text_from_file(file_path):
        """Extract text content from various file types"""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.pdf':
                return FileProcessingService._extract_from_pdf(file_path)
            elif file_extension in ['.jpg', '.jpeg', '.png', '.gif']:
                return FileProcessingService._extract_from_image(file_path)
            elif file_extension == '.docx':
                return FileProcessingService._extract_from_docx(file_path)
            elif file_extension in ['.ppt', '.pptx']:
                return FileProcessingService._extract_from_pptx(file_path)
            elif file_extension == '.txt':
                return FileProcessingService._extract_from_txt(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
                
        except Exception as e:
            raise Exception(f"Error processing file: {str(e)}")
    
    @staticmethod
    def _extract_from_pdf(file_path):
        """Extract text from PDF file"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
        
        return text.strip()
    
    @staticmethod
    def _extract_from_image(file_path):
        """Extract text from image using OCR (placeholder - would need OCR library like Tesseract)"""
        # For now, return a placeholder message
        # In production, you would use pytesseract or similar OCR library
        return f"Image file detected: {os.path.basename(file_path)}. OCR text extraction would be implemented here."
    
    @staticmethod
    def _extract_from_docx(file_path):
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            raise Exception(f"Error reading DOCX: {str(e)}")
        
        return text.strip()
    
    @staticmethod
    def _extract_from_pptx(file_path):
        """Extract text from PowerPoint file"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise Exception(f"Error reading PPTX: {str(e)}")
        
        return text.strip()
    
    @staticmethod
    def _extract_from_txt(file_path):
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            raise Exception(f"Error reading text file: {str(e)}")
    
    @staticmethod
    def get_file_info(file_path):
        """Get file information"""
        file_size = os.path.getsize(file_path)
        file_name = os.path.basename(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()
        
        return {
            'name': file_name,
            'size': file_size,
            'extension': file_extension,
            'type': FileProcessingService._get_file_type(file_extension)
        }
    
    @staticmethod
    def _get_file_type(extension):
        """Determine file type from extension"""
        if extension == '.pdf':
            return 'PDF Document'
        elif extension in ['.jpg', '.jpeg', '.png', '.gif']:
            return 'Image'
        elif extension == '.docx':
            return 'Word Document'
        elif extension in ['.ppt', '.pptx']:
            return 'PowerPoint Presentation'
        elif extension == '.txt':
            return 'Text File'
        else:
            return 'Unknown'
